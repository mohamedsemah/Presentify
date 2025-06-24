import io
import os
from typing import Dict, List
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches
from PIL import Image as PILImage


class ExportService:
    @staticmethod
    def export_to_pdf(presentation_data: Dict) -> bytes:
        """Export presentation to PDF"""
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        # Title page
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30,
            alignment=1  # Center alignment
        )

        story.append(Paragraph(presentation_data['title'], title_style))
        if presentation_data.get('description'):
            story.append(Paragraph(presentation_data['description'], styles['Normal']))
        story.append(Spacer(1, 0.5 * inch))

        # Slides
        for i, slide in enumerate(presentation_data['slides'], 1):
            # Slide title
            slide_title = f"Slide {i}: {slide.get('title', 'Untitled')}"
            story.append(Paragraph(slide_title, styles['Heading1']))

            # Slide content blocks
            for block in slide.get('content_blocks', []):
                if block['type'] == 'text':
                    story.append(Paragraph(block['content'], styles['Normal']))
                elif block['type'] == 'image' and block.get('file_path'):
                    try:
                        img = Image(block['file_path'], width=4 * inch, height=3 * inch)
                        story.append(img)
                    except:
                        story.append(Paragraph(f"[Image: {block.get('content', 'Image')}]", styles['Italic']))

                story.append(Spacer(1, 12))

            story.append(Spacer(1, 0.3 * inch))

        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()

    @staticmethod
    def export_to_pptx(presentation_data: Dict) -> bytes:
        """Export presentation to PowerPoint"""
        prs = PPTXPresentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = presentation_data['title']
        if presentation_data.get('description'):
            subtitle.text = presentation_data['description']

        # Content slides
        for slide_data in presentation_data['slides']:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)

            title = slide.shapes.title
            title.text = slide_data.get('title', 'Untitled Slide')

            # Add content blocks
            content_top = Inches(2)
            for block in slide_data.get('content_blocks', []):
                if block['type'] == 'text':
                    left = Inches(1)
                    width = Inches(8)
                    height = Inches(1)

                    textbox = slide.shapes.add_textbox(left, content_top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = block['content']
                    content_top += height + Inches(0.2)

                elif block['type'] == 'image' and block.get('file_path'):
                    try:
                        left = Inches(2)
                        width = Inches(6)
                        height = Inches(4)
                        slide.shapes.add_picture(block['file_path'], left, content_top, width, height)
                        content_top += height + Inches(0.2)
                    except:
                        pass

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
